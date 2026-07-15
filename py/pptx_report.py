"""Generador del reporte PPTX de productividad (porta R/utils_crear_pptx.R)."""

import math

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from py.charts import fecha_larga_es, mes_anio_es
from py.utils_comunes import calcular_fecha_corte

# ---------------------------------------------------------------------------
# Colores
# ---------------------------------------------------------------------------
COL_ROJO_CHILLON = "FF0000"
COL_AMARILLO_CHILLON = "FFC107"
COL_VERDE_CHILLON = "00B050"
COL_MUTED = "6B7280"
COL_BORDE = "D1D5DB"
COL_TEXTO = "111827"
COL_VERDE = "1E5B4F"
COL_GUINDA = "611232"
COL_DORADO = "A57F2C"

_TIPO_PROC_PPTX_MAP = {
    "consulta total": "total_consultas",
    "general": "consulta_gral",
    "especialidad": "consulta_esp",
    "qx": "qx",
    "egresos": "egresos",
}

MAPA_TITULOS_CONSULTAS = {
    "total_consultas": "Consultas totales",
    "consulta_gral": "Consulta general",
    "consulta_esp": "Especialidad",
    "qx": "Procedimientos quirúrgicos",
    "egresos": "Egresos",
}

MAPA_TITULOS_CURP = {
    "total_consultas": "Consultas totales",
    "consulta_gral": "Consulta general",
    "consulta_esp": "Especialidad",
    "qx": "Intervenidas",
    "egresos": "Egresadas",
}


def _rgb(hexcolor: str) -> RGBColor:
    return RGBColor.from_string(hexcolor.lstrip("#"))


# ---------------------------------------------------------------------------
# Helpers de placeholders (equivalente a officer::ph_location_label)
# ---------------------------------------------------------------------------

def _layout_idx_map(layout) -> dict:
    return {ph.name: ph.placeholder_format.idx for ph in layout.placeholders}


def find_placeholder(layout, slide, name):
    idx_map = _layout_idx_map(layout)
    if name not in idx_map:
        return None
    return slide.placeholders[idx_map[name]]


def _bbox(ph):
    return ph.left, ph.top, ph.width, ph.height


def _remove(shape):
    shape._element.getparent().remove(shape._element)


def add_slide(prs: Presentation, master, layout_name: str):
    layout = next(l for l in master.slide_layouts if l.name == layout_name)
    slide = prs.slides.add_slide(layout)
    return slide, layout


def set_placeholder_text(layout, slide, name: str, text: str):
    ph = find_placeholder(layout, slide, name)
    if ph is None:
        return
    p = ph.text_frame.paragraphs[0]
    if p.runs:
        run = p.runs[0]
        for extra in list(p.runs[1:]):
            extra._r.getparent().remove(extra._r)
    else:
        run = p.add_run()
    run.text = str(text)


def place_native_chart(slide, layout, name: str, funcion_dibujo, *args, **kwargs):
    """Resuelve el placeholder `name` y llama a
    `funcion_dibujo(slide, box, *args, **kwargs)` para dibujar formas nativas
    (editables en PowerPoint) dentro de su geometría, en vez de pegar una
    imagen (como hacía la vieja place_picture)."""
    ph = find_placeholder(layout, slide, name)
    if ph is None:
        return
    box = _bbox(ph)
    _remove(ph)
    funcion_dibujo(slide, box, *args, **kwargs)


def _set_cell_border(cell, color: str = COL_MUTED, width_pt: float = 1.0):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    line_w = str(int(Pt(width_pt)))
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tcPr.find(qn(tag))
        if existing is not None:
            tcPr.remove(existing)
        ln = tcPr.makeelement(qn(tag), {"w": line_w, "cap": "flat"})
        solid_fill = ln.makeelement(qn("a:solidFill"), {})
        srgb = solid_fill.makeelement(qn("a:srgbClr"), {"val": color})
        solid_fill.append(srgb)
        ln.append(solid_fill)
        tcPr.append(ln)


def place_table(
    slide, layout, name: str, df: pd.DataFrame,
    w: list | None = None,
    header_negro: str = "3B3B3B",
    menta: str = "D9F2EE",
    size_header: int = 10,
    size_body: int = 9,
):
    ph = find_placeholder(layout, slide, name)
    if ph is None:
        return
    left, top, width, height = _bbox(ph)
    _remove(ph)

    n_rows, n_cols = df.shape[0] + 1, df.shape[1]
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    for j, colname in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = "" if colname == "indicador" else str(colname)
        if colname == "indicador":
            cell.text_frame.paragraphs[0].add_run().text = "Indicador"
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(header_negro)
        _set_cell_border(cell)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(size_header)
                run.font.color.rgb = _rgb("FFFFFF")

    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            cell = table.cell(i + 1, j)
            valor = df.iat[i, j]
            cell.text = "" if pd.isna(valor) else str(valor)
            _set_cell_border(cell)
            if j == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(menta)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("FFFFFF")
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                if not p.runs:
                    p.add_run().text = ""
                for run in p.runs:
                    run.font.size = Pt(size_body)

    if w:
        total_in = sum(w)
        for j, width_in in enumerate(w):
            table.columns[j].width = Emu(int(width * (width_in / total_in)))

    return table


# ---------------------------------------------------------------------------
# Formas nativas de bajo nivel para gráficas (editables en PowerPoint)
# ---------------------------------------------------------------------------
MESES_EN_ABBR = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
                  "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]


def _forma_rect(slide, left, top, width, height, color_hex, transparencia_pct=0):
    """Rectángulo nativo relleno. `transparencia_pct`: 0 (opaco) a 100 (invisible)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, int(left), int(top), int(max(width, 1)), int(max(height, 1)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    if transparencia_pct:
        srgb = shp.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int((100 - transparencia_pct) * 1000))})
        srgb.append(alpha)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _forma_texto(slide, left, top, width, height, texto, size=9, color_hex=COL_TEXTO,
                 bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 rotation=0, wrap=True):
    tb = slide.shapes.add_textbox(int(left), int(top), int(max(width, 1)), int(max(height, 1)))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, linea in enumerate(str(texto).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = linea
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = _rgb(color_hex)
        r.font.name = "Calibri"
    if rotation:
        tb.rotation = rotation
    return tb


def _forma_ovalo(slide, cx, cy, radio, color_hex):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, int(cx - radio), int(cy - radio), int(radio * 2), int(radio * 2))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _forma_linea_quebrada(slide, puntos, color_hex, width_pt=1.5):
    """Línea poligonal nativa y editable (freeform) que conecta `puntos`
    (lista de (x_emu, y_emu)) — a diferencia de una imagen, en PowerPoint se
    puede mover cada punto con doble clic."""
    if len(puntos) < 2:
        return None
    fb = slide.shapes.build_freeform(int(puntos[0][0]), int(puntos[0][1]), scale=1.0)
    fb.add_line_segments([(int(x), int(y)) for x, y in puntos[1:]], close=False)
    shp = fb.convert_to_shape()
    shp.fill.background()
    shp.line.color.rgb = _rgb(color_hex)
    shp.line.width = Pt(width_pt)
    shp.shadow.inherit = False
    return shp


def _forma_flecha_vertical(slide, x, y_top, y_bottom, color_hex, width_pt=1.2):
    """Conector vertical con punta de flecha arriba (en `y_top`)."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, int(x), int(y_top), int(x), int(y_bottom))
    conn.line.color.rgb = _rgb(color_hex)
    conn.line.width = Pt(width_pt)
    ln = conn.line._get_or_add_ln()
    head = ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(head)
    return conn


def _escala_bonita(valor_max, n_pasos_objetivo=5):
    """Devuelve (paso, valor_redondeado) con incrementos 'bonitos' (1/2/2.5/5/10 x 10^n),
    para que el eje Y muestre números redondos (12,000,000) en vez de un valor
    arbitrario (11,776,152)."""
    if valor_max <= 0:
        return 1, n_pasos_objetivo
    bruto = valor_max / n_pasos_objetivo
    exponente = math.floor(math.log10(bruto))
    base = 10 ** exponente
    paso = base * 10
    for m in (1, 2, 2.5, 5, 10):
        if bruto <= m * base:
            paso = m * base
            break
    return paso, math.ceil(valor_max / paso) * paso


def _dibujar_eje_y(slide, plot_l, plot_w, baseline, plot_h, ymax_eje, valores_marca,
                   ancho_etiqueta=None, gridlines=False):
    """Dibuja las marcas del eje Y (texto, en `valores_marca`) y opcionalmente
    líneas de rejilla, mapeadas proporcionalmente contra `ymax_eje`."""
    ancho_etiqueta = ancho_etiqueta or Emu(900000)
    for valor in valores_marca:
        frac = (valor / ymax_eje) if ymax_eje else 0
        y = baseline - frac * plot_h
        _forma_texto(
            slide, plot_l - ancho_etiqueta - Pt(4), y - Pt(7), ancho_etiqueta, Pt(14),
            fmt_num(valor), size=8, color_hex=COL_MUTED, align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        if gridlines and valor > 0:
            _forma_rect(slide, plot_l, y, plot_w, Pt(0.6), "E5E7EB")


# ---------------------------------------------------------------------------
# Gráficas de barras (histórico 2020-2025 y 2024-2026) — formas nativas
# ---------------------------------------------------------------------------
def dibujar_grafica_barras(slide, box, categorias, totales, avances, titulo,
                           colores_total, colores_avance,
                           etiquetas_total=None, etiquetas_avance=None,
                           title_size=15):
    """Barras de 'total' (claro) con 'avance' superpuesto (oscuro), con
    etiquetas de valor — equivalente nativo de grafica_planeacion_*."""
    L, T, W, H = box
    n = len(categorias)
    if n == 0:
        return

    if etiquetas_total is None:
        etiquetas_total = [fmt_num(t) for t in totales]
    if etiquetas_avance is None:
        etiquetas_avance = [fmt_num(a) for a in avances]

    alto_titulo = int(H * 0.13)
    alto_categoria = Pt(16)
    margen_izq = int(W * 0.12)
    espacio_etiqueta_sup = int(H * 0.16)

    _forma_texto(slide, L, T, W, alto_titulo, titulo, size=title_size,
                color_hex=COL_MUTED, bold=True, align=PP_ALIGN.CENTER)

    plot_l = L + margen_izq
    plot_t = T + alto_titulo + espacio_etiqueta_sup
    plot_w = W - margen_izq
    plot_h = H - alto_titulo - espacio_etiqueta_sup - alto_categoria
    baseline = plot_t + plot_h

    ymax_datos = max(max(totales, default=0), max(avances, default=0), 1)
    paso, ymax_redondeado = _escala_bonita(ymax_datos, 5)
    ymax_eje = max(ymax_redondeado, ymax_datos) * 1.16

    valores_marca = [i * paso for i in range(int(ymax_eje // paso) + 1)]
    _dibujar_eje_y(slide, plot_l, plot_w, baseline, plot_h, ymax_eje, valores_marca,
                  ancho_etiqueta=margen_izq - Pt(4))

    slot_w = plot_w / n
    bar_w = slot_w * 0.60

    for i, cat in enumerate(categorias):
        slot_l = plot_l + i * slot_w
        bar_l = slot_l + (slot_w - bar_w) / 2

        h_total = (totales[i] / ymax_eje) * plot_h if ymax_eje else 0
        h_avance = (avances[i] / ymax_eje) * plot_h if ymax_eje else 0

        if h_total > 0:
            _forma_rect(slide, bar_l, baseline - h_total, bar_w, h_total, colores_total[i])
        if h_avance > 0:
            _forma_rect(slide, bar_l, baseline - h_avance, bar_w, h_avance, colores_avance[i])

        n_lineas_tot = etiquetas_total[i].count("\n") + 1
        alto_et = Pt(13) * n_lineas_tot
        _forma_texto(slide, slot_l, baseline - h_total - alto_et - Pt(3), slot_w, alto_et,
                    etiquetas_total[i], size=9.5, color_hex="000000", bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

        if h_avance > Pt(24):
            n_lineas_av = etiquetas_avance[i].count("\n") + 1
            alto_ea = Pt(12) * n_lineas_av
            _forma_texto(slide, slot_l, baseline - h_avance + Pt(3), slot_w, alto_ea,
                        etiquetas_avance[i], size=8.5, color_hex="FFFFFF", bold=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

        _forma_texto(slide, slot_l, baseline + Pt(3), slot_w, alto_categoria,
                    str(cat), size=11, color_hex=COL_MUTED, bold=True, align=PP_ALIGN.CENTER)


def dibujar_grafica_planeacion_historica(slide, box, df, col_total, col_avance, titulo,
                                         beige="D9D2BE", verde="2F6F63"):
    """Barras 2020-2025: total (beige) con avance (verde) superpuesto."""
    anios = list(range(2020, 2026))
    d = df.copy()
    d["anio_num"] = d["anio"].astype(int)
    d = d[d["anio_num"].isin(anios)].set_index("anio_num").reindex(anios)
    totales = d[col_total].fillna(0).astype(float).tolist()
    avances = d[col_avance].fillna(0).astype(float).tolist()

    dibujar_grafica_barras(
        slide, box, categorias=[str(a) for a in anios],
        totales=totales, avances=avances, titulo=titulo,
        colores_total=[beige] * len(anios), colores_avance=[verde] * len(anios),
    )


def dibujar_grafica_planeacion_2024_2026(slide, box, df, col_total, col_avance, titulo,
                                         beige="D9D2BE", verde="2F6F63",
                                         beige_2026="A99F86", verde_2026="1E5B4F"):
    """Barras 2024-2026 con etiquetas especiales de 'Meta 2026' / 'Avance'."""
    anios = [2024, 2025, 2026]
    d = df.copy()
    d["anio_num"] = d["anio"].astype(int)
    d = d[d["anio_num"].isin(anios)].set_index("anio_num").reindex(anios)
    totales = d[col_total].fillna(0).astype(float).tolist()
    avances = d[col_avance].fillna(0).astype(float).tolist()

    etiquetas_total, etiquetas_avance = [], []
    for a, tot, av in zip(anios, totales, avances):
        if a == 2026:
            etiquetas_total.append(f"Meta 2026\n{fmt_num(tot)}")
            pct = (av / tot) if tot > 0 else np.nan
            pct_txt = "s/d" if pd.isna(pct) else f"{int(round(pct * 100))}%"
            etiquetas_avance.append(f"Avance\n{fmt_num(av)}\n({pct_txt})")
        else:
            etiquetas_total.append(fmt_num(tot))
            etiquetas_avance.append(fmt_num(av))

    colores_total = [beige_2026 if a == 2026 else beige for a in anios]
    colores_avance = [verde_2026 if a == 2026 else verde for a in anios]

    dibujar_grafica_barras(
        slide, box, categorias=[str(a) for a in anios],
        totales=totales, avances=avances, titulo=titulo,
        colores_total=colores_total, colores_avance=colores_avance,
        etiquetas_total=etiquetas_total, etiquetas_avance=etiquetas_avance,
    )


# ---------------------------------------------------------------------------
# Gráfica de serie temporal (consultas/procedimientos por mes) — formas nativas
# ---------------------------------------------------------------------------
def dibujar_grafica_consultas_periodos(slide, box, df, fecha_inicio="2022-08-01", fecha_fin=None,
                                       titulo="Consultas totales del IMSS Bienestar",
                                       color_linea="6B6B6B", verde_punto="1F5B50",
                                       fill_2223="EFEFEF", fill_2024="E9DDCC",
                                       fill_2025="F4F0EA", fill_2026="E9DDCC",
                                       fill_valuebox="B99C6D"):
    """Serie temporal mensual con bandas por periodo y anotaciones, dibujada
    con formas nativas de PowerPoint (línea freeform editable, rectángulos,
    óvalos y textos) en vez de una imagen."""
    L, T, W, H = box
    hoy = pd.Timestamp.today().normalize()
    inicio_mes_actual = hoy.replace(day=1)
    if fecha_fin is None:
        fecha_fin = inicio_mes_actual
    else:
        fecha_fin = pd.Timestamp(fecha_fin)
    fecha_inicio = pd.Timestamp(fecha_inicio)

    d = df.copy()
    d["fecha"] = pd.to_datetime(d["fecha"])
    d = d[
        (d["fecha"] >= fecha_inicio) & (d["fecha"] <= fecha_fin) & (d["fecha"] < inicio_mes_actual)
    ].sort_values("fecha")

    # Título (dentro del área del gráfico, arriba del todo)
    alto_titulo = Pt(20)
    _forma_texto(slide, L, T, W, alto_titulo, titulo, size=15, color_hex=COL_MUTED,
                bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    if d.empty:
        _forma_texto(slide, L, T + alto_titulo, W, H - alto_titulo,
                    "Sin datos suficientes para este período", size=12, color_hex=COL_MUTED)
        return

    ymax = float(d["consultas_totales"].max())
    ymin = float(d["consultas_totales"].min())
    ymax_eje = max(ymax * 1.48, 1)

    fecha_fin_banda = (fecha_fin + pd.offsets.MonthBegin(1)).normalize()
    bandas = [
        (fecha_inicio, pd.Timestamp("2024-01-01"), fill_2223, "2022–2023\nAños de transición"),
        (pd.Timestamp("2024-01-01"), pd.Timestamp("2025-01-01"), fill_2024, "2024\nPrimer año de operación"),
        (pd.Timestamp("2025-01-01"), pd.Timestamp("2026-01-01"), fill_2025, "2025\nSegundo año de operación"),
        (pd.Timestamp("2026-01-01"), fecha_fin_banda, fill_2026, "2026\nTercer año de operación"),
    ]

    dias_pad = max((fecha_fin_banda - fecha_inicio).days * 0.035, 10)
    x0_ord = fecha_inicio.toordinal() - dias_pad
    x1_ord = fecha_fin_banda.toordinal() + dias_pad

    margen_izq = int(W * 0.075)
    margen_der = int(W * 0.015)
    margen_sup = alto_titulo + int(H * 0.05)
    margen_inf = int(H * 0.17)

    plot_l = L + margen_izq
    plot_t = T + margen_sup
    plot_w = W - margen_izq - margen_der
    plot_h = H - margen_sup - margen_inf
    baseline = plot_t + plot_h

    def _ord(valor_fecha):
        return pd.Timestamp(valor_fecha).toordinal()

    def xmap(valor_fecha):
        frac = (_ord(valor_fecha) - x0_ord) / (x1_ord - x0_ord)
        return plot_l + frac * plot_w

    def ymap(valor):
        frac = (valor / ymax_eje) if ymax_eje else 0
        return baseline - frac * plot_h

    _forma_rect(slide, L, plot_t, W, plot_h, "FFFFFF")

    for xmin_b, xmax_b, fill, _lab in bandas:
        x_l, x_r = xmap(xmin_b), xmap(xmax_b)
        _forma_rect(slide, x_l, plot_t, x_r - x_l, plot_h, fill)

    _paso_eje, _ = _escala_bonita(ymax, 5)
    valores_marca_y = [i * _paso_eje for i in range(int(ymax_eje // _paso_eje) + 1)]
    _dibujar_eje_y(slide, plot_l, plot_w, baseline, plot_h, ymax_eje, valores_marca_y,
                  ancho_etiqueta=margen_izq - Pt(4), gridlines=True)

    ultimos_3 = d.tail(3)
    if not ultimos_3.empty:
        xmin_sub = ultimos_3["fecha"].min() - pd.Timedelta(days=15)
        xmax_sub = ultimos_3["fecha"].max() + pd.Timedelta(days=15)
        x_l, x_r = xmap(xmin_sub), xmap(xmax_sub)
        _forma_rect(slide, x_l, plot_t, x_r - x_l, plot_h, "B22222", transparencia_pct=82)
        _forma_texto(
            slide, xmap(d["fecha"].max()) - Pt(70), ymap(ymax * 1.08), Pt(140), Pt(26),
            "Posible subregistro\ntemporal", size=8.5, color_hex="7A1E3A", bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    puntos_linea = [(xmap(f), ymap(v)) for f, v in zip(d["fecha"], d["consultas_totales"])]
    _forma_linea_quebrada(slide, puntos_linea, color_linea, width_pt=1.3)
    for x, y in puntos_linea:
        _forma_ovalo(slide, x, y, Pt(1.6), color_linea)

    mes_destacado = fecha_fin.month
    puntos_destacados = d[(d["fecha"].dt.month == mes_destacado) & (d["fecha"].dt.year < 2026)]
    for _, fila in puntos_destacados.iterrows():
        x, y = xmap(fila["fecha"]), ymap(fila["consultas_totales"])
        _forma_ovalo(slide, x, y, Pt(4.5), verde_punto)
        etiqueta = f"{fmt_num(fila['consultas_totales'])}\n{fila['fecha'].strftime('%b-%Y').title()}"
        _forma_texto(slide, x - Pt(45), y - Pt(38), Pt(90), Pt(30), etiqueta,
                    size=8, color_hex="000000", bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

    for xmin_b, xmax_b, _fill, lab in bandas:
        centro = xmap(xmin_b) + (xmap(xmax_b) - xmap(xmin_b)) / 2
        _forma_texto(slide, centro - Pt(70), ymap(ymax * 1.32), Pt(140), Pt(26), lab,
                    size=8, color_hex="000000", bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    x_decreto = xmap(pd.Timestamp("2022-08-15"))
    _forma_flecha_vertical(slide, x_decreto, ymap(ymax * 1.02), ymap(ymin * 0.95), verde_punto)
    _forma_texto(slide, x_decreto + Pt(6), ymap(ymax * 1.05) - Pt(4), Pt(120), Pt(26),
                "Decreto de creación\ndel IMSS Bienestar", size=7.5, color_hex="000000",
                bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    fecha_ultimo_valor = d["fecha"].max()
    valor_ultimo = d.loc[d["fecha"] == fecha_ultimo_valor, "consultas_totales"].iloc[0]
    x_vb, y_vb = xmap(fecha_ultimo_valor), ymap(valor_ultimo)
    vb_w, vb_h = Pt(72), Pt(30)
    vb = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, int(x_vb - vb_w / 2), int(y_vb - vb_h - Pt(10)),
        int(vb_w), int(vb_h))
    vb.fill.solid()
    vb.fill.fore_color.rgb = _rgb(fill_valuebox)
    vb.line.fill.background()
    vb.shadow.inherit = False
    tf = vb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lineas_valuebox = [fmt_num(valor_ultimo), fecha_ultimo_valor.strftime("%b %Y").title()]
    for i, linea in enumerate(lineas_valuebox):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = linea
        r.font.size = Pt(9.5)
        r.font.bold = True
        r.font.color.rgb = _rgb("FFFFFF")
        r.font.name = "Calibri"

    eje_inicio = fecha_inicio.replace(day=1) - pd.DateOffset(months=2)
    eje_fin = fecha_fin_banda + pd.DateOffset(months=2)
    for tick in pd.date_range(eje_inicio, eje_fin, freq="2MS"):
        x = xmap(tick)
        if x < plot_l - Pt(5) or x > plot_l + plot_w + Pt(5):
            continue
        etiqueta = f"{MESES_EN_ABBR[tick.month - 1]}-{str(tick.year)[2:]}"
        _forma_texto(slide, x - Pt(20), baseline + Pt(4), Pt(40), Pt(22), etiqueta,
                    size=7.5, color_hex=COL_MUTED, bold=False, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.TOP, rotation=45, wrap=False)


# ---------------------------------------------------------------------------
# Tarjetas institucionales (equivalente a rvg::dml + grid, con formas nativas)
# ---------------------------------------------------------------------------

def fmt_delta(x) -> dict:
    if x is None or (isinstance(x, float) and np.isnan(x)):
        return {"label": "s/d", "col": COL_MUTED, "icon": ""}
    if x > 0:
        return {"label": f"+{x:.0f}%", "col": COL_VERDE, "icon": "▲ "}
    if x < 0:
        return {"label": f"{x:.0f}%", "col": COL_GUINDA, "icon": "▼ "}
    return {"label": "0%", "col": COL_MUTED, "icon": "• "}


def fmt_num(x) -> str:
    return f"{x:,.0f}"


def elige_acento(
    var_2025, var_2024,
    verde: str = COL_VERDE_CHILLON, amarillo: str = COL_AMARILLO_CHILLON, rojo: str = COL_ROJO_CHILLON,
) -> str:
    v25_neg = pd.notna(var_2025) and var_2025 < 0
    v24_neg = pd.notna(var_2024) and var_2024 < 0
    if v25_neg and v24_neg:
        return rojo
    if v25_neg != v24_neg:
        return amarillo
    return verde


def _card_textbox(slide, left, top, width, height, x_frac, y_top_frac, h_frac, parts, size, align=PP_ALIGN.LEFT):
    box_left = left + int(width * x_frac)
    box_top = top + int(height * y_top_frac)
    box_width = width - int(width * x_frac) - int(width * 0.04)
    box_height = int(height * h_frac)
    tb = slide.shapes.add_textbox(box_left, box_top, box_width, box_height)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    for text, color, bold in parts:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Calibri"
        run.font.color.rgb = _rgb(color)
    return tb


def dibujar_card(
    slide, left, top, width, height,
    numero, titulo, var_vs_2025, var_vs_2024,
    acento: str = COL_VERDE, size_num: float = 30, size_titulo: float = 13, size_delta: float = 10.5,
):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb("FFFFFF")
    card.line.color.rgb = _rgb(COL_BORDE)
    card.line.width = Pt(1)
    card.shadow.inherit = False
    card.text_frame.paragraphs[0].text = ""
    card.adjustments[0] = 0.08

    barra = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + int(width * 0.02), top + int(height * 0.05),
        int(width * 0.02), int(height * 0.90),
    )
    barra.fill.solid()
    barra.fill.fore_color.rgb = _rgb(acento)
    barra.line.fill.background()
    barra.shadow.inherit = False

    _card_textbox(slide, left, top, width, height, 0.08, 0.03, 0.36, [(fmt_num(numero), COL_DORADO, True)], size_num)
    _card_textbox(slide, left, top, width, height, 0.08, 0.36, 0.22, [(str(titulo), COL_TEXTO, True)], size_titulo)

    d25 = fmt_delta(var_vs_2025)
    d24 = fmt_delta(var_vs_2024)

    _card_textbox(
        slide, left, top, width, height, 0.08, 0.60, 0.18,
        [(f"{d25['icon']}vs 2025 ", COL_MUTED, False), (d25["label"], d25["col"], True)],
        size_delta,
    )
    _card_textbox(
        slide, left, top, width, height, 0.08, 0.78, 0.18,
        [(f"{d24['icon']}vs 2024 ", COL_MUTED, False), (d24["label"], d24["col"], True)],
        size_delta,
    )


def _agregar_valuebox(slide, layout, ph_name, df, metrica, titulo):
    ph = find_placeholder(layout, slide, ph_name)
    if ph is None:
        return
    left, top, width, height = _bbox(ph)
    _remove(ph)

    numero = valor_anio_col(df, metrica, 2026)
    numero = 0 if pd.isna(numero) else numero

    var_2025 = valor_anio_col(df, f"var_2026_vs_2025_{metrica}", 2026)
    var_2024 = valor_anio_col(df, f"var_2026_vs_2024_{metrica}", 2026)

    acento = elige_acento(var_2025, var_2024)

    dibujar_card(
        slide, left, top, width, height,
        numero=numero, titulo=titulo,
        var_vs_2025=var_2025, var_vs_2024=var_2024,
        acento=acento,
    )


def agregar_valueboxes(slide, layout, metricas_vb, datos_consulta_funcion, datos_curps):
    for i, metrica in enumerate(metricas_vb, start=1):
        _agregar_valuebox(slide, layout, f"arriba {i}", datos_consulta_funcion, metrica, MAPA_TITULOS_CONSULTAS[metrica])
        _agregar_valuebox(slide, layout, f"abajo {i}", datos_curps, metrica, MAPA_TITULOS_CURP[metrica])


# ---------------------------------------------------------------------------
# Helpers de datos (equivalentes 1:1 a utils_crear_pptx.R)
# ---------------------------------------------------------------------------

def obtener_col(df: pd.DataFrame, nombre: str, default=0) -> pd.Series:
    if nombre in df.columns:
        return df[nombre]
    return pd.Series([default] * len(df), index=df.index)


def rellenar_anios(df: pd.DataFrame, anios) -> pd.DataFrame:
    base = pd.DataFrame({"anio": list(anios)})
    out = base.merge(df, on="anio", how="left")
    numeric_cols = out.select_dtypes(include="number").columns
    out[numeric_cols] = out[numeric_cols].fillna(0)
    return out


def asegurar_columnas(df: pd.DataFrame, cols) -> pd.DataFrame:
    df = df.copy()
    for c in cols:
        if c not in df.columns:
            df[c] = 0
    return df


def valor_anio_col(df: pd.DataFrame, columna: str, anio_objetivo):
    if columna not in df.columns:
        return np.nan
    val = df.loc[df["anio"] == anio_objetivo, columna]
    if val.empty:
        return np.nan
    return val.iloc[0]


def hay_indicador_2026(df: pd.DataFrame, columna: str) -> bool:
    val = valor_anio_col(df, columna, 2026)
    return bool(pd.notna(val) and val != 0)


def tiene_dato_2026(df: pd.DataFrame, columna: str) -> bool:
    if columna not in df.columns:
        return False
    valores = df.loc[df["anio"] == 2026, columna]
    if valores.empty:
        return False
    return bool(((valores.notna()) & (valores != 0)).any())


def definir_layout_valueboxes(datos_consulta_funcion: pd.DataFrame) -> dict:
    hay_general = tiene_dato_2026(datos_consulta_funcion, "consulta_gral")
    hay_esp = tiene_dato_2026(datos_consulta_funcion, "consulta_esp")
    hay_qx = tiene_dato_2026(datos_consulta_funcion, "qx")
    hay_egresos = tiene_dato_2026(datos_consulta_funcion, "egresos")

    if hay_general and not hay_esp and not hay_qx and not hay_egresos:
        return {"layout": "2_valueboxes", "metricas": ["consulta_gral"]}
    if hay_general and hay_esp and not hay_qx and not hay_egresos:
        return {"layout": "6_valueboxes", "metricas": ["total_consultas", "consulta_gral", "consulta_esp"]}
    if hay_general and hay_esp and hay_qx and not hay_egresos:
        return {"layout": "8_valueboxes", "metricas": ["total_consultas", "consulta_gral", "consulta_esp", "qx"]}
    if hay_general and hay_esp and hay_qx and hay_egresos:
        return {
            "layout": "10_valueboxes",
            "metricas": ["total_consultas", "consulta_gral", "consulta_esp", "qx", "egresos"],
        }

    metricas_presentes = []
    if hay_general:
        metricas_presentes.append("consulta_gral")
    if hay_esp:
        metricas_presentes.append("consulta_esp")
    if hay_qx:
        metricas_presentes.append("qx")
    if hay_egresos:
        metricas_presentes.append("egresos")
    if len(metricas_presentes) >= 2:
        metricas_presentes = ["total_consultas"] + metricas_presentes

    n = len(metricas_presentes)
    if n <= 1:
        layout_fallback = "2_valueboxes"
    elif n == 3:
        layout_fallback = "6_valueboxes"
    elif n == 4:
        layout_fallback = "8_valueboxes"
    else:
        layout_fallback = "10_valueboxes"

    return {"layout": layout_fallback, "metricas": metricas_presentes}


def armar_tabla_dinamica(df: pd.DataFrame, indicadores: list, etiquetas: list, mes_nombre: str) -> pd.DataFrame:
    filas = []
    col_25 = f"{mes_nombre} 2025"
    col_26 = f"{mes_nombre} 2026"
    for ind, etq in zip(indicadores, etiquetas):
        v25 = valor_anio_col(df, ind, 2025)
        v26 = valor_anio_col(df, ind, 2026)

        if pd.isna(v25) or v25 == 0:
            crecimiento_str = "s/d"
        else:
            crecimiento = round((v26 - v25) / v25 * 100)
            crecimiento_str = f"+{crecimiento} %" if crecimiento > 0 else f"{crecimiento} %"

        v25_str = "s/d" if pd.isna(v25) else f"{v25:,.0f}"
        v26_str = "s/d" if pd.isna(v26) else f"{v26:,.0f}"

        filas.append({"indicador": etq, col_25: v25_str, col_26: v26_str, "Crecimiento anual": crecimiento_str})

    return pd.DataFrame(filas, columns=["indicador", col_25, col_26, "Crecimiento anual"])


# ---------------------------------------------------------------------------
# Función principal
# ---------------------------------------------------------------------------

def crear_reporte_productividad(
    codigo_clues: str,
    clues_info: pd.DataFrame,
    metas: pd.DataFrame,
    historicos: pd.DataFrame,
    procedimientos_personas: pd.DataFrame,
    ruta_master: str,
) -> Presentation:
    fecha_corte = calcular_fecha_corte()
    fecha_portada = fecha_larga_es(fecha_corte)
    fecha_corte_15 = pd.Timestamp(year=fecha_corte.year, month=fecha_corte.month, day=15)
    fecha_fin_graf = fecha_corte_15

    historicos = historicos.copy()
    historicos["fecha"] = pd.to_datetime(historicos["fecha"])
    historicos["anio"] = historicos["fecha"].dt.year

    procedimientos_personas = procedimientos_personas.copy()
    procedimientos_personas["fecha"] = procedimientos_personas["fecha"].astype(float).astype(int)
    procedimientos_personas["tipo_procedimiento"] = (
        procedimientos_personas["tipo_procedimiento"].map(_TIPO_PROC_PPTX_MAP).fillna(procedimientos_personas["tipo_procedimiento"])
    )

    hay_fila_2026 = bool((procedimientos_personas["fecha"] == 2026).any())

    clues_info_filtrado = clues_info[clues_info["clues_imb"] == codigo_clues]
    if clues_info_filtrado.empty:
        raise ValueError("No se encontró el codigo_clues en clues_info.")

    metas_filtrado = metas[metas["clues_imb"] == codigo_clues]
    if metas_filtrado.empty:
        metas_filtrado = pd.DataFrame([{
            "clues_imb": codigo_clues,
            "meta_general_anual": 0, "meta_especialidad_anual": 0,
            "meta_cirugia_anual": 0, "meta_egresos_anual": 0,
        }])

    meta_total_consultas = (
        metas_filtrado["meta_general_anual"].fillna(0).sum() + metas_filtrado["meta_especialidad_anual"].fillna(0).sum()
    )
    meta_qx = metas_filtrado["meta_cirugia_anual"].fillna(0).sum()

    prs = Presentation(ruta_master)
    master = prs.slide_masters[0]

    # Portada ------------------------------------------------------------
    slide, layout = add_slide(prs, master, "Portada 3")
    nombre_unidad = str(clues_info_filtrado["nombre_de_la_unidad"].iloc[0]).title()
    set_placeholder_text(
        layout, slide, "Título 1",
        f"Reporte de productividad médica\n{nombre_unidad} ({clues_info_filtrado['clues_imb'].iloc[0]})",
    )
    set_placeholder_text(layout, slide, "Marcador de contenido 2", fecha_portada)

    # Datos base -----------------------------------------------------------
    cols_metricas = ["consulta_gral", "consulta_esp", "qx", "total_consultas", "egresos"]

    historicos_agg = asegurar_columnas(
        historicos, ["consulta_general", "consulta_especialidad", "procedimientos_qx", "consulta_total", "egresos"]
    )
    datos_anual = historicos_agg.groupby("anio", as_index=False).agg(
        consulta_gral_anual=("consulta_general", "sum"),
        consulta_esp_anual=("consulta_especialidad", "sum"),
        qx_anual=("procedimientos_qx", "sum"),
        total_consultas_anual=("consulta_total", "sum"),
        egresos_anual=("egresos", "sum"),
    )
    datos_anual = rellenar_anios(datos_anual, range(2020, 2027))

    datos_consulta_funcion = (
        procedimientos_personas.rename(columns={"fecha": "anio"})
        .pivot_table(index="anio", columns="tipo_procedimiento", values="procedimientos", aggfunc="sum")
        .reset_index()
    )
    datos_consulta_funcion = rellenar_anios(datos_consulta_funcion, range(2024, 2027))
    datos_consulta_funcion = asegurar_columnas(datos_consulta_funcion, cols_metricas)
    datos_consulta_funcion = datos_consulta_funcion.merge(datos_anual, on="anio", how="left")
    numeric_cols = datos_consulta_funcion.select_dtypes(include="number").columns
    datos_consulta_funcion[numeric_cols] = datos_consulta_funcion[numeric_cols].fillna(0)

    datos_consulta_funcion["total_consultas_meta"] = np.where(
        datos_consulta_funcion["anio"] == 2026, meta_total_consultas, datos_consulta_funcion["total_consultas_anual"],
    )
    datos_consulta_funcion["qx_meta"] = np.where(
        datos_consulta_funcion["anio"] == 2026, meta_qx, datos_consulta_funcion["qx_anual"],
    )

    datos_curps = (
        procedimientos_personas.rename(columns={"fecha": "anio"})
        .pivot_table(index="anio", columns="tipo_procedimiento", values="personas", aggfunc="sum")
        .reset_index()
    )
    datos_curps = rellenar_anios(datos_curps, range(2024, 2027))
    datos_curps = asegurar_columnas(datos_curps, cols_metricas)
    numeric_cols = datos_curps.select_dtypes(include="number").columns
    datos_curps[numeric_cols] = datos_curps[numeric_cols].fillna(0)

    for col in cols_metricas:
        for ref in (2024, 2025):
            nombre_var = f"var_2026_vs_{ref}_{col}"
            for df_ in (datos_consulta_funcion, datos_curps):
                valor_2026 = valor_anio_col(df_, col, 2026)
                valor_ref = valor_anio_col(df_, col, ref)
                # R permite Inf/-Inf en este cálculo sin truronar (valor_2026==0 con
                # valor_ref!=0); en Python hay que evitarlo explícitamente para no
                # crashear con OverflowError al redondear un infinito.
                if pd.notna(valor_2026) and pd.notna(valor_ref) and valor_ref != 0 and valor_2026 != 0:
                    computed = round(100 * (1 - (valor_ref / valor_2026)))
                else:
                    computed = 0
                df_[nombre_var] = np.where(df_["anio"] == 2026, computed, 0)

    # Value boxes ------------------------------------------------------------
    config_vb = definir_layout_valueboxes(datos_consulta_funcion)
    layout_vb, metricas_vb = config_vb["layout"], config_vb["metricas"]

    slide, layout = add_slide(prs, master, layout_vb)
    set_placeholder_text(layout, slide, "Título 1", "Productividad IMSS Bienestar")
    set_placeholder_text(layout, slide, "fecha", f"Del 01 de enero al {fecha_portada}")
    agregar_valueboxes(slide, layout, metricas_vb, datos_consulta_funcion, datos_curps)

    # Diapo 3: 2024-2026, solo si hay 2026 ------------------------------------
    if hay_fila_2026:
        d = datos_consulta_funcion.copy()
        d["anio_num"] = d["anio"].astype(int)
        for c in ["consulta_gral", "consulta_esp", "qx", "egresos", "consulta_gral_anual", "consulta_esp_anual", "qx_anual"]:
            d[c] = d[c].fillna(0)

        d["total_consultas"] = np.select(
            [(d["consulta_gral"] > 0) & (d["consulta_esp"] > 0), d["consulta_gral"] > 0, d["consulta_esp"] > 0],
            [d["consulta_gral"] + d["consulta_esp"], d["consulta_gral"], d["consulta_esp"]],
            default=0,
        )

        d["total_consultas_meta"] = np.select(
            [
                (d["anio_num"] == 2026) & d["total_consultas_meta"].notna(),
                d["anio_num"] == 2026,
                (d["consulta_gral_anual"] > 0) & (d["consulta_esp_anual"] > 0),
                d["consulta_gral_anual"] > 0,
                d["consulta_esp_anual"] > 0,
                d["total_consultas_anual"] > 0,
            ],
            [
                d["total_consultas_meta"], d["total_consultas"],
                d["consulta_gral_anual"] + d["consulta_esp_anual"], d["consulta_gral_anual"],
                d["consulta_esp_anual"], d["total_consultas_anual"],
            ],
            default=d["total_consultas"],
        )

        d["qx_meta"] = np.select(
            [(d["anio_num"] == 2026) & d["qx_meta"].notna(), d["anio_num"] == 2026, d["qx_anual"] > 0],
            [d["qx_meta"], d["qx"], d["qx_anual"]],
            default=d["qx"],
        )

        datos_2024_2026 = d[d["anio_num"].isin([2024, 2025, 2026])]

        hay_consultas_2024_2026 = bool((datos_2024_2026["total_consultas"] > 0).any())
        hay_qx_2024_2026 = bool(((datos_2024_2026["qx"] > 0) | (datos_2024_2026["egresos"] > 0)).any())

        if hay_consultas_2024_2026:
            indicadores_consulta, etiquetas_consulta = [], []
            if hay_indicador_2026(datos_2024_2026, "consulta_gral"):
                indicadores_consulta.append("consulta_gral")
                etiquetas_consulta.append("Consultas generales")
            if hay_indicador_2026(datos_2024_2026, "consulta_esp"):
                indicadores_consulta.append("consulta_esp")
                etiquetas_consulta.append("Consultas de especialidad*")

            tabla_consultas = armar_tabla_dinamica(datos_2024_2026, indicadores_consulta, etiquetas_consulta, "Acumulado")

            if hay_qx_2024_2026:
                indicadores_proc, etiquetas_proc = [], []
                if hay_indicador_2026(datos_2024_2026, "qx"):
                    indicadores_proc.append("qx")
                    etiquetas_proc.append("Procedimientos quirúrgicos")
                if hay_indicador_2026(datos_2024_2026, "egresos"):
                    indicadores_proc.append("egresos")
                    etiquetas_proc.append("Egresos")

                tabla_proc = armar_tabla_dinamica(datos_2024_2026, indicadores_proc, etiquetas_proc, "Acumulado")

                slide, layout = add_slide(prs, master, "Historico consultas y procedimientos")
                set_placeholder_text(layout, slide, "Título 1", "Productividad IMSS Bienestar")
                place_native_chart(slide, layout, "Grafica 1", dibujar_grafica_planeacion_2024_2026,
                                   datos_2024_2026, "total_consultas_meta", "total_consultas", "Consultas totales")
                place_native_chart(slide, layout, "Grafica 2", dibujar_grafica_planeacion_2024_2026,
                                   datos_2024_2026, "qx_meta", "qx", "Procedimientos quirúrgicos")
                place_table(slide, layout, "tabla_1", tabla_consultas, w=[2.70, 0.90, 0.90, 0.80], size_header=8, size_body=8)
                place_table(slide, layout, "tabla_2", tabla_proc, w=[2.70, 0.90, 0.90, 0.80], size_header=8, size_body=8)
                set_placeholder_text(layout, slide, "fecha", f"Del 01 de enero al {fecha_portada}")
            else:
                slide, layout = add_slide(prs, master, "Historico consultas")
                set_placeholder_text(layout, slide, "Título 1", "Productividad IMSS Bienestar")
                place_native_chart(slide, layout, "Grafica 1", dibujar_grafica_planeacion_2024_2026,
                                   datos_2024_2026, "total_consultas_meta", "total_consultas", "Consultas totales")
                place_table(slide, layout, "tabla_1", tabla_consultas, w=[4.60, 1.35, 1.35, 1.40], size_header=11, size_body=10)
                set_placeholder_text(layout, slide, "fecha", f"Del 01 de enero al {fecha_portada}")

    # Diapo 4: históricos 2020-2025 -------------------------------------------
    dh = historicos.copy()
    dh["consulta_general_tmp"] = obtener_col(dh, "consulta_general") + obtener_col(dh, "consulta_gral")
    dh["consulta_esp_tmp"] = obtener_col(dh, "consulta_especialidad") + obtener_col(dh, "consulta_esp")
    dh["procedimientos_qx_tmp"] = obtener_col(dh, "procedimientos_qx")
    dh["egresos_tmp"] = obtener_col(dh, "egresos")

    mask = (
        (dh["fecha"].dt.year >= 2020)
        & (dh["fecha"].dt.year <= 2025)
        & (dh["fecha"].dt.strftime("%m-%d") <= fecha_corte_15.strftime("%m-%d"))
    )
    dh = dh[mask].copy()
    dh["anio"] = dh["fecha"].dt.year.astype(str)

    resumen = dh.groupby("anio").agg(
        consulta_gral=("consulta_general_tmp", "sum"),
        consulta_esp=("consulta_esp_tmp", "sum"),
        qx=("procedimientos_qx_tmp", "sum"),
        egresos=("egresos_tmp", "sum"),
    ).reset_index()

    resumen["total_consultas"] = np.select(
        [(resumen["consulta_gral"] > 0) & (resumen["consulta_esp"] > 0), resumen["consulta_gral"] > 0, resumen["consulta_esp"] > 0],
        [resumen["consulta_gral"] + resumen["consulta_esp"], resumen["consulta_gral"], resumen["consulta_esp"]],
        default=0,
    )

    datos_anual_str = datos_anual.copy()
    datos_anual_str["anio"] = datos_anual_str["anio"].astype(str)
    resumen = resumen.merge(datos_anual_str, on="anio", how="left")
    for c in ["consulta_gral_anual", "consulta_esp_anual", "qx_anual"]:
        resumen[c] = resumen[c].fillna(0)

    resumen["total_consultas_anual"] = np.select(
        [
            (resumen["consulta_gral_anual"] > 0) & (resumen["consulta_esp_anual"] > 0),
            resumen["consulta_gral_anual"] > 0,
            resumen["consulta_esp_anual"] > 0,
        ],
        [resumen["consulta_gral_anual"] + resumen["consulta_esp_anual"], resumen["consulta_gral_anual"], resumen["consulta_esp_anual"]],
        default=0,
    )
    datos_historicos_2020_2025 = resumen.sort_values("anio")

    hay_consultas_2020 = bool((datos_historicos_2020_2025["total_consultas"] > 0).any())
    hay_qx_2020 = bool(
        ((datos_historicos_2020_2025["qx"] > 0) | (datos_historicos_2020_2025["egresos"] > 0)).any()
    )

    if hay_consultas_2020:
        if hay_qx_2020:
            slide, layout = add_slide(prs, master, "1_Historico consultas y procedimientos")
            set_placeholder_text(layout, slide, "Título 1", "Productividad IMSS Bienestar")
            place_native_chart(slide, layout, "Grafica 1", dibujar_grafica_planeacion_historica,
                               datos_historicos_2020_2025, "total_consultas_anual", "total_consultas",
                               "Consultas totales")
            place_native_chart(slide, layout, "Grafica 2", dibujar_grafica_planeacion_historica,
                               datos_historicos_2020_2025, "qx_anual", "qx", "Procedimientos quirúrgicos")
            set_placeholder_text(layout, slide, "fecha", f"Del 01 de enero al {fecha_portada}")
        else:
            slide, layout = add_slide(prs, master, "1_Historico consultas")
            set_placeholder_text(layout, slide, "Título 1", "Productividad IMSS Bienestar")
            place_native_chart(slide, layout, "Grafica 1", dibujar_grafica_planeacion_historica,
                               datos_historicos_2020_2025, "total_consultas_anual", "total_consultas",
                               "Consultas totales")
            set_placeholder_text(layout, slide, "fecha", f"Del 01 de enero al {fecha_portada}")

    # Diapo 5 ------------------------------------------------------------
    serie_mensual_consultas = historicos.copy()
    serie_mensual_consultas["fecha"] = serie_mensual_consultas["fecha"].values.astype("datetime64[M]")
    serie_mensual_consultas = serie_mensual_consultas.dropna(subset=["fecha", "consulta_total"])
    serie_mensual_consultas = (
        serie_mensual_consultas.groupby("fecha", as_index=False)["consulta_total"]
        .sum()
        .rename(columns={"consulta_total": "consultas_totales"})
        .sort_values("fecha")
    )

    if len(serie_mensual_consultas) > 0:
        slide, layout = add_slide(prs, master, "Una grafica")
        set_placeholder_text(layout, slide, "Título 1", "Consultas totales por mes (2022-2026)")
        place_native_chart(
            slide, layout, "ft", dibujar_grafica_consultas_periodos,
            serie_mensual_consultas, fecha_inicio="2022-08-01", fecha_fin=str(fecha_fin_graf.date()),
            titulo=f"Consultas totales del IMSS Bienestar (agosto 2022 – {mes_anio_es(fecha_fin_graf)})",
        )

    # Diapo 6 ------------------------------------------------------------
    if hay_fila_2026 and hay_indicador_2026(datos_consulta_funcion, "qx"):
        serie_mensual_pq = historicos.copy()
        serie_mensual_pq["fecha"] = serie_mensual_pq["fecha"].values.astype("datetime64[M]")
        serie_mensual_pq = serie_mensual_pq.dropna(subset=["fecha", "procedimientos_qx"])
        serie_mensual_pq = (
            serie_mensual_pq.groupby("fecha", as_index=False)["procedimientos_qx"]
            .sum()
            .rename(columns={"procedimientos_qx": "consultas_totales"})
            .sort_values("fecha")
        )

        if len(serie_mensual_pq) > 0:
            slide, layout = add_slide(prs, master, "Una grafica")
            set_placeholder_text(layout, slide, "Título 1", "Procedimientos quirúrgicos por mes (2022-2026)")
            place_native_chart(
                slide, layout, "ft", dibujar_grafica_consultas_periodos,
                serie_mensual_pq, fecha_inicio="2022-08-01", fecha_fin=str(fecha_fin_graf.date()),
                titulo=f"Procedimientos quirúrgicos del IMSS Bienestar (agosto 2022 – {mes_anio_es(fecha_fin_graf)})",
            )

    return prs
